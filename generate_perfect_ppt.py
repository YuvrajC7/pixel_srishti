# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK_SLIDE_LAYOUT = 6

DARK_BLUE = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
ACCENT = RGBColor(14, 165, 233)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(71, 85, 105)

def add_header(slide, text):
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = DARK_BLUE
    return title

# --- SLIDE 1 ---
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK_BLUE

title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "PIXEL-Srishti"
p.font.bold = True
p.font.size = Pt(72)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
tf = subtitle.text_frame
p = tf.paragraphs[0]
p.text = "Agentic Geospatial AI for Watershed Monitoring"
p.font.size = Pt(28)
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER

team = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1))
tf = team.text_frame
p = tf.paragraphs[0]
p.text = "Team: JALDI THE LATE | SIH 2026 Review 1"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# --- SLIDE 2 ---
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
add_header(slide, "The Escalating Geospatial Bottleneck")

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Despite massive growth in satellite imagery, ground-level watershed development is stalled by slow manual GIS processing."
p.font.size = Pt(20)
p.font.color.rgb = GRAY

p = tf.add_paragraph()
p.text = "\n- Massive Data Sprawl: The Satellite Data Market is doubling, reaching $14.75B by 2034."
p.font.size = Pt(18)
p.font.color.rgb = DARK_BLUE

p = tf.add_paragraph()
p.text = "\n- The Skill Gap: District officers lack the advanced QGIS/ArcGIS training required to fuse Optical + SAR datasets."
p.font.size = Pt(18)

p = tf.add_paragraph()
p.text = "\n- Operational Latency: Identifying crop health or water body changes manually takes 8-40 hours per district. We need instant answers."
p.font.size = Pt(18)

chart_data = CategoryChartData()
chart_data.categories = ['2020', '2022', '2025', '2030', '2034']
chart_data.add_series('Satellite Data Market (USD Billion)', (2.8, 3.5, 5.72, 10.2, 14.75))
chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.8), Inches(1.5), Inches(6), Inches(4.5), chart_data).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

# --- SLIDE 3 ---
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
add_header(slide, "System Architecture: Agentic Orchestration")

def draw_box(slide, x, y, w, h, text, color, font_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = DARK_BLUE
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = font_color
    p.alignment = PP_ALIGN.CENTER
    return shape

draw_box(slide, 1, 3.5, 2, 1, "Field Admin\n(Natural Language)", GRAY, WHITE)
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.2), Inches(3.8), Inches(0.6), Inches(0.4)).fill.solid()
draw_box(slide, 4, 3, 2.5, 2, "Agentic Orchestrator\n(LangGraph LLM)\n\nIntent Parsing &\nTool Selection", BLUE, WHITE)
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.7), Inches(2.2), Inches(0.6), Inches(0.3)).fill.solid()
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.7), Inches(3.8), Inches(0.6), Inches(0.3)).fill.solid()
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.7), Inches(5.4), Inches(0.6), Inches(0.3)).fill.solid()
draw_box(slide, 7.5, 1.8, 2.5, 1, "Optical-SAR Fusion\n(Change Detection)", ACCENT, WHITE)
draw_box(slide, 7.5, 3.4, 2.5, 1, "Thematic Segmentation\n(U-Net)", ACCENT, WHITE)
draw_box(slide, 7.5, 5.0, 2.5, 1, "VQA / Captioning\n(LLaVA)", ACCENT, WHITE)
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(10.2), Inches(3.8), Inches(0.6), Inches(0.4)).fill.solid()
draw_box(slide, 11, 3, 2, 2, "React GUI\n(Interactive Map\n& GeoJSON Masks)", DARK_BLUE, WHITE)

# --- SLIDE 4 ---
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
add_header(slide, "Real-World Impact of Watershed AI")

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "By automating GIS workflows, PIXEL-Srishti directly enables rapid watershed intervention, leading to measurable rural impact."
p.font.size = Pt(20)
p.font.color.rgb = GRAY

p = tf.add_paragraph()
p.text = "\n- 75% Increase in Crop Productivity: Rapid detection of soil moisture and drainage allows optimized interventions (Source: Kadwanchi Study)."
p.font.size = Pt(18)
p.font.color.rgb = DARK_BLUE

p = tf.add_paragraph()
p.text = "\n- Reduced Migration: Effective watershed tracking reduces seasonal drought migration by 20-50%."
p.font.size = Pt(18)

chart_data = CategoryChartData()
chart_data.categories = ['Pre-Intervention', 'Post-Intervention']
chart_data.add_series('Crop Productivity Index', (100, 175.4))
chart_data.add_series('Groundwater Level (m) x10', (32, 51.7))
chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(6.5), Inches(1.5), Inches(6.5), Inches(5), chart_data).chart
chart.has_legend = True

prs.save('PIXEL_Srishti_Master_Review1.pptx')