# -*- coding: utf-8 -*-
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK_SLIDE_LAYOUT = 6

# Colors
DARK_BG = RGBColor(15, 23, 42)
LIGHT_BG = RGBColor(248, 250, 252)
ACCENT_BLUE = RGBColor(37, 99, 235)
ACCENT_LIGHT = RGBColor(56, 189, 248)
TEXT_DARK = RGBColor(30, 41, 59)
TEXT_LIGHT = RGBColor(241, 245, 249)
GRAY = RGBColor(71, 85, 105)

def create_title(slide, text, color=TEXT_DARK):
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = color
    return title

def create_body(slide, text, left=0.5, top=1.5, width=12.3, height=5.5, color=TEXT_DARK, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for line in text.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
    return box

# --- SLIDE 1: TITLE ---
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = DARK_BG

box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
p = box.text_frame.paragraphs[0]
p.text = "PIXEL-Srishti"
p.font.bold = True
p.font.size = Pt(72)
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
p = box.text_frame.paragraphs[0]
p.text = "An Agentic Geospatial AI Assistant for Multi-Sensor Satellite Image Analysis\nand Watershed Development Monitoring"
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_LIGHT
p.alignment = PP_ALIGN.CENTER

box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.5))
p = box.text_frame.paragraphs[0]
p.text = "Team: JALDI THE LATE\nMinistry of Rural Development & ISRO–SAC"
p.font.size = Pt(18)
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

# --- SLIDE 2: PROBLEM STATEMENT ---
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT_BG
create_title(slide, "1. The Problem Statement")
content = (
    "Currently, satellite image analysis tools are highly fragmented and built for single tasks:\n\n"
    "- Siloed Capabilities: One tool classifies land types, another spots changes, another handles VQA.\n"
    "- Technical Barrier: Combining Optical, Multispectral, and SAR data requires specialized GIS training.\n"
    "- The Gap: Organizations like ISRO collect massive amounts of high-quality data (SRISHTI-DRISHTI), but "
    "it sits unused because field teams lack automated, natural-language tools to convert raw data into insights.\n\n"
    "Consequence: Decision-makers lack fast, clear answers about land use, drainage, and water body changes."
)
create_body(slide, content)

# --- SLIDE 3: CONVERGENCE RATIONALE ---
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT_BG
create_title(slide, "2. Convergence Rationale")
content = (
    "The core gap is the absence of an intelligent, query-driven interpretation layer between raw satellite "
    "imagery and field-level decisions.\n\n"
    "PIXEL-Srishti introduces a general-purpose Agentic Remote-Sensing Assistant. While watershed monitoring "
    "is the flagship real-world application, the underlying architecture builds the reusable, multi-task AI "
    "backbone required by the SatQuery AI mandate.\n\n"
    "Keywords: Agentic VLM, Multimodal Analysis, Optical-SAR Pairs, VQA, GIS-Free Decision Making."
)
create_body(slide, content)

# --- SLIDE 4: SOLUTION OBJECTIVES ---
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT_BG
create_title(slide, "3. Solution Objectives")
content = (
    "- Agentic Vision-Language System: Answers queries over single images, optical-SAR pairs, and bi-temporal series.\n"
    "- Domain Grounding: Fine-tuning on BigEarthNet, VRSBench, RSVQA, and CDVQA for accurate, auditable outputs.\n"
    "- Automated Thematic Mapping: Instantly generates land-use, drainage, and vegetation maps from 30m resolution data.\n"
    "- Spatial Change Detection: Identifies new water-conservation structures and vegetation recovery over time.\n"
    "- Orchestration Layer: Automatically selects the correct specialist ML model based on the user's plain-English query."
)
create_body(slide, content)

# --- SLIDE 5: SYSTEM ARCHITECTURE (BLOCK DIAGRAM) ---
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = DARK_BG
create_title(slide, "4. System Architecture (4-Layer Framework)", TEXT_LIGHT)

def draw_box(slide, x, y, w, h, text, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.color.rgb = TEXT_LIGHT
    tf = shape.text_frame; p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    return shape

draw_box(slide, 1, 2, 2.5, 4.5, "1. Data Ingestion\n\nGeoTIFF/TIFF\nOptical, SAR,\nMultispectral\n\nSRISHTI-DRISHTI Tiles", GRAY)
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.6), Inches(4), Inches(0.4), Inches(0.3)).fill.solid()

draw_box(slide, 4.1, 2, 2.5, 4.5, "2. Agentic Controller\n\nLLM-based Router\nLangGraph pattern\n\nInterprets intent\nChecks modality\nRoutes to models", ACCENT_BLUE)
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.7), Inches(4), Inches(0.4), Inches(0.3)).fill.solid()

draw_box(slide, 7.2, 1.5, 3, 1, "Remote-Sensing VQA / Caption", ACCENT_LIGHT)
draw_box(slide, 7.2, 2.7, 3, 1, "Text-Guided Grounding Model", ACCENT_LIGHT)
draw_box(slide, 7.2, 3.9, 3, 1, "Bi-Temporal Change Model", ACCENT_LIGHT)
draw_box(slide, 7.2, 5.1, 3, 1, "Watershed Segmentation (U-Net)", ACCENT_LIGHT)
draw_box(slide, 7.2, 6.3, 3, 1, "Optical-SAR Fusion Model", ACCENT_LIGHT)

slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(10.3), Inches(4), Inches(0.4), Inches(0.3)).fill.solid()
draw_box(slide, 10.8, 2, 2, 4.5, "4. Presentation Layer\n\nReact Web GUI\nInteractive Maps\nChange Overlays\nExecution Traces", GRAY)

# --- SLIDE 6: EFFICIENCY COMPARISON (BAR CHART) ---
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT_BG
create_title(slide, "5. Impact: GIS Processing Time Reduction")

content = (
    "PIXEL-Srishti fundamentally shifts the bottleneck from manual execution to automated inference.\n\n"
    "By allowing field planners to use natural language, tasks that previously required data downloading, "
    "co-registration, and manual GIS mapping are executed instantly."
)
create_body(slide, content, left=0.5, top=1.5, width=6, height=5)

chart_data = CategoryChartData()
chart_data.categories = ['Data Ingestion', 'Feature Extraction', 'Change Detection', 'Report Generation']
chart_data.add_series('Traditional GIS Workflow (Hours)', (4.0, 12.0, 8.0, 3.0))
chart_data.add_series('PIXEL-Srishti Automation (Hours)', (0.1, 0.5, 0.3, 0.1))
chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(6.8), Inches(1.5), Inches(6), Inches(5), chart_data).chart
chart.has_legend = True

# --- SLIDE 7: RISKS & MITIGATION ---
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT_BG
create_title(slide, "6. Feasibility & Mitigation")

content = (
    "1. Risk: Limited GPU time for fine-tuning multiple models.\n"
    "   Mitigation: Use parameter-efficient fine-tuning (LoRA/QLoRA) on pre-trained backbones.\n\n"
    "2. Risk: Scarcity of labelled watershed-specific ground truth.\n"
    "   Mitigation: Bootstrap labels from existing geo-tagged field photographs and SRISHTI-DRISHTI thematic layers.\n\n"
    "3. Risk: Agentic controller mis-routing queries.\n"
    "   Mitigation: Constrain tool registry with strict input-compatibility checks and log auditable traces.\n\n"
    "4. Risk: Cloud cover affecting optical imagery.\n"
    "   Mitigation: SAR fusion pathway provides all-weather, day-and-night fallback."
)
create_body(slide, content, font_size=18)

# --- SLIDE 8: EXPECTED OUTCOMES ---
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT_BG
create_title(slide, "7. Expected Outcomes")

content = (
    "- Instant Actionable Insights: A working web application replacing weeks of manual interpretation.\n\n"
    "- Automated Thematic Maps: Direct extraction of land use, drainage, and vegetation from 30m satellite data.\n\n"
    "- Reusable Backbone: The agentic framework satisfies the SatQuery AI mandate and can be redeployed for disaster management and urban planning.\n\n"
    "- Transparency & Auditing: Every answer ships with an auditable execution trace (models used, confidence scores), suitable for government administrative use."
)
create_body(slide, content)

# --- SLIDE 9: BUDGET & FUTURE SCOPE ---
slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT_BG
create_title(slide, "8. Budget Estimate & Future Scope")

# Budget Pie Chart based on PDF
chart_data = ChartData()
chart_data.categories = ['GPU Compute (Cloud)', 'Hardware Peripherals', 'Expert Consultation', 'Satellite Data', 'Miscellaneous', 'Storage', 'Software APIs']
chart_data.add_series('Budget (INR)', (45000, 15000, 12000, 10000, 10000, 8000, 5000))
chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5), chart_data).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM

content = (
    "Future Scope:\n"
    "- Extend framework to crop-yield estimation and flood mapping without redesigning the controller.\n"
    "- Integrate with state/district Watershed Management Information Systems.\n"
    "- Mobile-first interface for field engineers with limited connectivity.\n"
    "- Incorporate IoT soil-moisture feeds for data fusion."
)
create_body(slide, content, left=6.5, top=1.5, width=6, height=5, font_size=18)

prs.save('PIXEL_Srishti_Final_Review1.pptx')