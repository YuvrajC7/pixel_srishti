from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Layouts
TITLE_SLIDE_LAYOUT = 0
BULLET_SLIDE_LAYOUT = 1
SECTION_HEADER_LAYOUT = 2
BLANK_SLIDE_LAYOUT = 6

# --- SLIDE 1: TITLE ---
slide_layout = prs.slide_layouts[TITLE_SLIDE_LAYOUT]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "PIXEL-Srishti"
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.size = Pt(54)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

subtitle.text = "Agentic Geospatial Assistant for Watershed Monitoring\n\nTeam: JALDI THE LATE\nSmart India Hackathon 2026 - Review 1"
subtitle.text_frame.paragraphs[0].font.size = Pt(24)

# --- SLIDE 2: THE PROBLEM ---
slide_layout = prs.slide_layouts[BULLET_SLIDE_LAYOUT]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "The Core Problem"
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

tf = body_shape.text_frame
tf.text = "Vast amounts of satellite data remain underutilized at the ground level:"
p = tf.add_paragraph()
p.text = "Data Silos: Optical, SAR, and multispectral data are processed separately."
p.level = 1
p = tf.add_paragraph()
p.text = "High Barrier to Entry: Processing raw GeoTIFFs requires complex GIS software (QGIS/ArcGIS) and specialized training."
p.level = 1
p = tf.add_paragraph()
p.text = "Slow Decision Making: Field administrators cannot instantly query data to track water body changes or vegetation recovery."
p.level = 1

# --- SLIDE 3: OUR SOLUTION ---
slide = prs.slides.add_slide(slide_layout)
title_shape = slide.shapes.title
body_shape = slide.shapes.placeholders[1]
title_shape.text = "Our Proposed Solution"
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

tf = body_shape.text_frame
tf.text = "A single, intelligent conversational interface that democratizes GIS:"
p = tf.add_paragraph()
p.text = "Natural Language Interface: Administrators simply ask questions in plain English (e.g., \"Show me the drainage lines here\")."
p.level = 1
p = tf.add_paragraph()
p.text = "Agentic Orchestration: The AI automatically selects the correct specialist model (Change Detection, Segmentation) to answer the query."
p.level = 1
p = tf.add_paragraph()
p.text = "Visual & Textual Output: Returns precise thematic map overlays (GeoJSON masks) directly onto an interactive dashboard."
p.level = 1

# --- SLIDE 4: SYSTEM ARCHITECTURE ---
slide = prs.slides.add_slide(slide_layout)
title_shape = slide.shapes.title
body_shape = slide.shapes.placeholders[1]
title_shape.text = "System Architecture (Block Flow)"
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

tf = body_shape.text_frame
tf.text = "How PIXEL-Srishti processes a field query:"
p = tf.add_paragraph()
p.text = "1. Data Ingestion Layer: Validates and pre-processes SRISHTI-DRISHTI tiles."
p.level = 1
p = tf.add_paragraph()
p.text = "2. Agentic Controller: Parses the natural language question using an LLM."
p.level = 1
p = tf.add_paragraph()
p.text = "3. Specialist Model Registry: Invokes the required ML pipeline (Optical-SAR Fusion, Siamese Change-Detection, Grounding)."
p.level = 1
p = tf.add_paragraph()
p.text = "4. Presentation Layer: Pushes actionable geo-coordinates and masks to the Web GUI."
p.level = 1

# --- SLIDE 5: TECHNOLOGY STACK ---
slide = prs.slides.add_slide(slide_layout)
title_shape = slide.shapes.title
body_shape = slide.shapes.placeholders[1]
title_shape.text = "Technology Stack"
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

tf = body_shape.text_frame
tf.text = "Built for performance and scalability:"
p = tf.add_paragraph()
p.text = "Frontend: React.js, React-Leaflet (Interactive Mapping), Tailwind CSS"
p.level = 1
p = tf.add_paragraph()
p.text = "Backend & Orchestration: FastAPI (Python), LangGraph/LangChain"
p.level = 1
p = tf.add_paragraph()
p.text = "Geospatial Data Processing: Rasterio, GeoPandas, GDAL"
p.level = 1
p = tf.add_paragraph()
p.text = "AI/ML Backbones: Grounding DINO, LLaVA-architecture vision models, Siamese U-Net"
p.level = 1

# --- SLIDE 6: PROJECT IMPACT ---
slide = prs.slides.add_slide(slide_layout)
title_shape = slide.shapes.title
body_shape = slide.shapes.placeholders[1]
title_shape.text = "Impact & Scalability"
title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

tf = body_shape.text_frame
tf.text = "Why this solution transforms watershed development:"
p = tf.add_paragraph()
p.text = "Replaces weeks of manual GIS processing with instant, automated insights."
p.level = 1
p = tf.add_paragraph()
p.text = "Eliminates the GIS training bottleneck for local administrators."
p.level = 1
p = tf.add_paragraph()
p.text = "Highly Scalable: The agentic framework can be easily expanded for crop-yield estimation, flood mapping, and urban sprawl."
p.level = 1

prs.save('PIXEL_Srishti_Review1.pptx')
